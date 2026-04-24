"""
Build slides from slides.yaml using U.S. Bank discussion template.

Uses Layout 12 ("Title and content 1") — white background with:
    PH 0  = Title    (0.5", 0.5") 9.0" x 1.2"
    PH 1  = Content  (0.5", 2.1") 9.0" x 4.9"  [OBJECT type 7]

And Layout 23 ("Title two subtitle and two content 3") for Slide 3:
    PH 0  = Title          (0.5", 0.5") 9.0" x 1.2"
    PH 10 = Left subtitle  (0.5", 1.9") 4.4" x 0.7"
    PH 1  = Left content   (0.5", 2.6") 4.4" x 3.9"
    PH 11 = Right subtitle (5.1", 1.9") 4.4" x 0.7"
    PH 2  = Right content  (5.1", 2.6") 4.4" x 3.9"

Usage:
    1. Place template as U_S_Bank_standard_discussion_temp_cleaned.pptx
    2. Run: python build_slides_template.py
    3. Open output in PowerPoint

To inspect layouts: python build_slides_template.py --inspect

Dependencies: pip install python-pptx pyyaml
"""

import yaml
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn


TEMPLATE_PATH = "U_S_Bank_standard_discussion_temp_cleaned.pptx"
YAML_PATH = "slides.yaml"
OUTPUT_PATH = "20260423_Agentic_AI_Validation_Best_Practices.pptx"

LAYOUT_CONTENT = 12       # "Title and content 1" — white, PH 0 + PH 1
LAYOUT_TWO_COL = 23       # "Title two subtitle and two content 3" — PH 0, 10, 1, 11, 2


# ---------------------------------------------------------------------------
# Inspect
# ---------------------------------------------------------------------------

def inspect_template(path):
    prs = Presentation(path)
    print(f"\nTemplate: {path}")
    print(f"Slide size: {prs.slide_width / 914400:.1f}\" x "
          f"{prs.slide_height / 914400:.1f}\"\n")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}: \"{layout.name}\"")
        for ph in layout.placeholders:
            print(f"    Placeholder {ph.placeholder_format.idx}: "
                  f"\"{ph.name}\" ({ph.placeholder_format.type}) "
                  f"@ ({ph.left/914400:.1f}\", {ph.top/914400:.1f}\") "
                  f"w={ph.width/914400:.1f}\" h={ph.height/914400:.1f}\"")
        print()


# ---------------------------------------------------------------------------
# Helper: populate a placeholder with structured paragraphs
# ---------------------------------------------------------------------------

def fill_placeholder(slide, ph_idx, paragraphs):
    """
    Fill a placeholder with a list of paragraph dicts.
    Each dict can have: text, bold, italic, size, space_after

    Handles OBJECT (type 7) placeholders by writing directly to
    paragraphs[0] to activate it, then cleaning up extra default
    paragraphs before adding new ones.
    """
    if ph_idx not in slide.placeholders:
        print(f"  Warning: placeholder {ph_idx} not found on slide, skipping")
        return

    ph = slide.placeholders[ph_idx]
    tf = ph.text_frame

    if not paragraphs:
        return

    # Write first paragraph to activate the placeholder
    first = paragraphs[0]
    p = tf.paragraphs[0]
    p.text = first.get("text", "")
    p.space_after = Pt(first.get("space_after", 2))
    if "size" in first:
        p.font.size = Pt(first["size"])
    if "bold" in first:
        p.font.bold = first["bold"]
    if "italic" in first:
        p.font.italic = first["italic"]

    # Remove any extra default paragraphs from the template
    p_elements = tf._txBody.findall(qn("a:p"))
    for extra_p in p_elements[1:]:
        tf._txBody.remove(extra_p)

    # Add remaining paragraphs
    for para in paragraphs[1:]:
        p = tf.add_paragraph()
        p.text = para.get("text", "")
        p.space_after = Pt(para.get("space_after", 2))
        if "size" in para:
            p.font.size = Pt(para["size"])
        if "bold" in para:
            p.font.bold = para["bold"]
        if "italic" in para:
            p.font.italic = para["italic"]


# ---------------------------------------------------------------------------
# Slide 1: Regulatory Foundation
#   Layout 12 — PH 0 (title), PH 1 (content body)
# ---------------------------------------------------------------------------

def build_slide_1(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    # Title
    slide.placeholders[0].text = data["title"]

    # Body: subtitle line, then quotes, then footer
    body = []

    # Subtitle as first line in body (since Layout 12 has no subtitle PH)
    body.append({"text": data["subtitle"], "italic": True, "size": 11,
                 "space_after": 8})

    for q in data["quotes"]:
        body.append({"text": q["label"], "bold": True, "size": 10,
                     "space_after": 0})
        body.append({"text": q["text"], "size": 10, "space_after": 0})
        body.append({"text": q["source"], "italic": True, "size": 9,
                     "space_after": 6})

    if data.get("footer"):
        body.append({"text": "", "size": 6, "space_after": 4})
        body.append({"text": data["footer"], "italic": True, "size": 9})

    fill_placeholder(slide, 1, body)


# ---------------------------------------------------------------------------
# Slide 2: Proposed Template Changes
#   Layout 12 — PH 0 (title), PH 1 (content body)
# ---------------------------------------------------------------------------

def build_slide_2(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    # Title
    slide.placeholders[0].text = data["title"]

    # Body: subtitle, then cards, then callout
    body = []

    body.append({"text": data["subtitle"], "italic": True, "size": 11,
                 "space_after": 8})

    for card in data["cards"]:
        body.append({"text": card["heading"], "bold": True, "size": 10,
                     "space_after": 0})
        body.append({"text": card["section"], "italic": True, "size": 8,
                     "space_after": 0})
        body.append({"text": card["body"], "size": 9, "space_after": 6})

    if data.get("callout"):
        body.append({"text": "", "size": 6, "space_after": 2})
        body.append({"text": "Core Principle", "bold": True, "size": 10,
                     "space_after": 0})
        body.append({"text": data["callout"], "italic": True, "size": 9})

    fill_placeholder(slide, 1, body)


# ---------------------------------------------------------------------------
# Slide 3: Best Practices (two-column)
#   Layout 23 — PH 0 (title), PH 10 (left sub), PH 1 (left body),
#               PH 11 (right sub), PH 2 (right body)
# ---------------------------------------------------------------------------

def build_slide_3(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO_COL])

    # PH 0: Title
    slide.placeholders[0].text = data["title"]

    # PH 10: Left subtitle
    slide.placeholders[10].text = data["left_header"]

    # PH 1: Left body
    left = []
    for sec in data["left_sections"]:
        left.append({"text": sec["heading"], "bold": True, "size": 9,
                     "space_after": 0})
        for bullet in sec["bullets"]:
            left.append({"text": f"- {bullet}", "size": 8, "space_after": 1})
        left.append({"text": "", "size": 4, "space_after": 3})
    fill_placeholder(slide, 1, left)

    # PH 11: Right subtitle
    slide.placeholders[11].text = data["right_header"]

    # PH 2: Right body
    right = []
    for sec in data["right_sections"]:
        right.append({"text": sec["heading"], "bold": True, "size": 9,
                      "space_after": 0})
        right.append({"text": sec["body"], "size": 8, "space_after": 4})
    fill_placeholder(slide, 2, right)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    if "--inspect" in sys.argv:
        inspect_template(TEMPLATE_PATH)
        return

    with open(YAML_PATH, "r") as f:
        content = yaml.safe_load(f)

    prs = Presentation(TEMPLATE_PATH)

    builders = {
        "title": build_slide_1,
        "cards": build_slide_2,
        "two_column": build_slide_3,
    }

    for slide_data in content["slides"]:
        layout_key = slide_data["layout"]
        builder = builders.get(layout_key)
        if builder:
            builder(prs, slide_data)
        else:
            print(f"Warning: unknown layout '{layout_key}', skipping.")

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
